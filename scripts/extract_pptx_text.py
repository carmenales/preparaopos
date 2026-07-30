"""
Extrae texto e imágenes de una presentación PPTX a Markdown.

Por qué un script aparte de extract_pdf_text.py: leer un PPTX convertido
a PDF (o su texto con herramientas de PDF) desordena el contenido, porque
cada cuadro de texto de la diapositiva es un objeto independiente sin
orden de lectura garantizado en el PDF resultante. Leyendo el .pptx
nativo con python-pptx se respeta el orden real de los shapes de cada
diapositiva.

Misma convención de salida que extract_pdf_text.py:
    <output_md>
    <output_md sin extensión>/images/diapositiva-N-img-M.<ext>
"""

import re
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _sanitize_filename(name: str) -> str:
    name = re.sub(r"[^a-zA-Z0-9._-]+", "-", name).strip("-")
    return name or "img"


def _shape_text(shape) -> str:
    if shape.has_text_frame:
        return shape.text_frame.text.strip()

    if shape.has_table:
        rows = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    return ""


def extract_pptx_to_markdown(input_pptx: Path, output_md: Path, extract_images: bool = True) -> None:
    if not input_pptx.exists():
        raise FileNotFoundError(f"No existe la presentación: {input_pptx}")

    output_md.parent.mkdir(parents=True, exist_ok=True)

    images_dir = output_md.parent / output_md.stem / "images"
    images_dir_relative = f"{output_md.stem}/images"

    presentation = Presentation(str(input_pptx))

    lines = [
        f"# Texto extraído: {input_pptx.name}",
        "",
        "> Texto extraído automáticamente desde una presentación PPTX. "
        "El orden de los cuadros de texto dentro de cada diapositiva puede no "
        "coincidir exactamente con el orden visual — revisa antes de dar por buena la estructura.",
        "",
        f"- Archivo origen: `{input_pptx.as_posix()}`",
        f"- Número de diapositivas: {len(presentation.slides)}",
        "",
    ]

    total_images = 0

    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"## Diapositiva {slide_index}")
        lines.append("")

        text_found = False
        image_index = 0

        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and extract_images:
                image_index += 1
                total_images += 1
                image = shape.image
                filename = _sanitize_filename(f"diapositiva-{slide_index}-img-{image_index}.{image.ext}")

                images_dir.mkdir(parents=True, exist_ok=True)
                (images_dir / filename).write_bytes(image.blob)

                lines.append(f"![Imagen de la diapositiva {slide_index}]({images_dir_relative}/{filename})")
                lines.append("")
                continue

            text = _shape_text(shape)
            if text:
                text_found = True
                lines.append(text)
                lines.append("")

        if not text_found and image_index == 0:
            lines.append("_Diapositiva sin texto ni imágenes reconocidas._")
            lines.append("")

        # Notas del orador: en apuntes de academia a veces llevan
        # explicaciones que no están en la propia diapositiva.
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append("**Notas del orador:**")
                lines.append("")
                lines.append(notes_text)
                lines.append("")

    output_md.write_text("\n".join(lines), encoding="utf-8")

    if total_images:
        print(f"Extraídas {total_images} imágenes en: {images_dir}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Extrae texto e imágenes de un PPTX a Markdown.")
    parser.add_argument("input_pptx", help="Ruta del PPTX de origen")
    parser.add_argument("output_md", help="Ruta del Markdown de salida")
    parser.add_argument(
        "--no-images",
        action="store_true",
        help="No extraer imágenes incrustadas, solo texto.",
    )

    args = parser.parse_args()

    extract_pptx_to_markdown(
        input_pptx=Path(args.input_pptx),
        output_md=Path(args.output_md),
        extract_images=not args.no_images,
    )


if __name__ == "__main__":
    main()
